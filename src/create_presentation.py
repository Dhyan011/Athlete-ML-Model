import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_pptx='PLAYHACK_ML_Presentation.pptx'):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Palette definition
    NAVY = RGBColor(15, 23, 42)        # #0F172A
    CARD_BG = RGBColor(30, 41, 59)     # #1E293B
    LIGHT_BG = RGBColor(248, 250, 252) # #F8FAFC
    DARK_TEXT = RGBColor(30, 41, 59)
    MUTED_TEXT = RGBColor(100, 116, 139)
    BLUE_ACCENT = RGBColor(37, 99, 235) # #2563EB
    EMERALD = RGBColor(16, 185, 129)
    WHITE = RGBColor(255, 255, 255)
    
    def add_header(slide, title_text, subtitle_text):
        # Header banner
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.color.rgb = BLUE_ACCENT
        p2.font.bold = True
        
    def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=RGBColor(226, 232, 240)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.color.rgb = NAVY
    
    # Accent line
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.5), Inches(1.2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE_ACCENT
    line.line.color.rgb = BLUE_ACCENT
    
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PLAYHACK — ML TRACK"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248)
    
    p = tf.add_paragraph()
    p.text = "Multi-Modal Athlete Injury Risk & Recovery Duration Forecasting"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p = tf.add_paragraph()
    p.text = "Predicting the onset timing, probability, and rehabilitation trajectory from wearable telemetry."
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(148, 163, 184)
    
    # Author box
    ab = s1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(10.9), Inches(1.0))
    tf_a = ab.text_frame
    pa = tf_a.paragraphs[0]
    pa.text = "IIT Guwahati Sports Board X Technical Board Hackathon | Submission Presentation"
    pa.font.size = Pt(13)
    pa.font.color.rgb = RGBColor(203, 213, 225)
    pa.font.bold = True

    # ==========================================
    # SLIDE 2: PROBLEM DEFINITION & FRAMEWORK
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. Challenge Overview & Problem Formulation", 
               "TEMPORAL SPLIT: 30-DAY OBSERVATION HORIZON → 30-DAY FUTURE RISK HORIZON")
    
    # Left Card: Two Windows Formulation
    add_card(s2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Dual Time-Window Architecture"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    bullets = [
        "Observation Window (Days 1–30): High-frequency telemetry (daily activity, sleep, hourly HR, intensities, training sessions, baseline biometrics).",
        "Risk Window (Days 31–60): Hidden future evaluation horizon.",
        "Zero-Leakage Constraint: Strictly prohibited from using signals past Day 30 for feature construction.",
        "Scale: 3,000 professional athletes across 6 sports (Football, Badminton, Athletics, Tennis, Basketball, Volleyball)."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        
    # Right Card: Three Prediction Targets
    add_card(s2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb = s2.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The Three Prediction Targets"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    targets = [
        ("Target A: Binary Injury Classification", "Predict whether athlete sustains an injury in Days 31–60 (0 or 1). Base incidence: 35.0%."),
        ("Target B1: Onset Day Offset", "For injured athletes, predict exact day in risk window (1 to 30) when injury occurs. Required for all rows."),
        ("Target B2: Recovery Duration", "Number of days (5 to 20) the athlete remains sidelined before returning to play. Required for all rows."),
        ("Evaluation Penalty Scheme", "Missed injuries (False Negatives) trigger a heavy penalty of 30 days applied to both timing predictions!")
    ]
    for title, desc in targets:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 3: EXPLORATORY DATA ANALYSIS (EDA)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. Exploratory Data Analysis & Domain Findings",
               "KEY DISCOVERY: PEAK WORKLOAD SPIKES (+44.8%) & CONTACT SPORT RECOVERY BIMODALITY")
    
    # Embed Figure 1
    if os.path.exists('figures/01_sport_and_recovery_analysis.png'):
        s3.shapes.add_picture('figures/01_sport_and_recovery_analysis.png', Inches(0.8), Inches(1.5), Inches(6.8), Inches(2.6))
        
    # Embed Figure 2
    if os.path.exists('figures/02_workload_spikes_and_acwr.png'):
        s3.shapes.add_picture('figures/02_workload_spikes_and_acwr.png', Inches(0.8), Inches(4.3), Inches(6.8), Inches(2.6))
        
    # Right Column: Insights Card
    add_card(s3, Inches(7.8), Inches(1.5), Inches(4.7), Inches(5.4))
    tb = s3.shapes.add_textbox(Inches(8.0), Inches(1.7), Inches(4.3), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Empirical Insights"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    insights = [
        "Workload Spike Footprint: Injured athletes had a mean single-day peak of 14,982 steps vs 10,344 for uninjured athletes (+44.84% spike).",
        "Acute Workload Surges (ACWR): Athletes crossing ACWR > 1.5 in the final observation week showed 3.2x higher injury incidence.",
        "Sport Bimodality in Recovery: Contact sports (Football: 14.1d, Basketball: 14.5d) require significantly longer rehab than court sports (Badminton: 9.9d, Tennis: 10.3d).",
        "Sleep Deficit Amplifiers: Injured athletes experienced acute drops in minimum nightly sleep (401m vs 413m) during heavy training blocks."
    ]
    for ins in insights:
        p = tf.add_paragraph()
        p.text = "• " + ins
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 4: FEATURE ENGINEERING ENGINE
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Domain-Specific Feature Engineering",
               "101 HIGH-RESOLUTION BIOMETRIC & SPORTS SCIENCE INDICATORS CONSTRUCTED")
    
    cards_data = [
        ("1. Workload & ACWR (Sports Science)", [
            "Acute (7-day) vs Chronic (30-day) step ratios",
            "Peak-to-mean workload spike deltas",
            "Very active & fairly active minutes surges",
            "Linear velocity slope of 30-day training volume"
        ]),
        ("2. Sleep & Autonomic Recovery", [
            "Sleep efficiency (Time Asleep / Time in Bed)",
            "Sleep Regularity Index (CV of sleep duration)",
            "Cumulative sleep debt (< 480 mins ideal)",
            "Acute sleep loss in final 7 observation days"
        ]),
        ("3. Cardiovascular & Exertion", [
            "Nocturnal resting HR (2 AM – 6 AM average & min)",
            "Peak exertion maximum heart rate",
            "Heart rate reserve (Peak HR - Nocturnal Min)",
            "Cardiac strain duration (> 140 bpm hours)"
        ]),
        ("4. Athlete Baselines & Interactions", [
            "Baseline BMI & body composition drift",
            "Experience ratio (Years Playing / Age)",
            "Sport type & contact risk categorization",
            "Prior season injury count risk multiplier"
        ])
    ]
    
    coords = [
        (Inches(0.8), Inches(1.5)),
        (Inches(6.8), Inches(1.5)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.8), Inches(4.3))
    ]
    
    for (title, feats), (left, top) in zip(cards_data, coords):
        add_card(s4, left, top, Inches(5.7), Inches(2.6))
        tb = s4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.3), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        for f in feats:
            p = tf.add_paragraph()
            p.text = "• " + f
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 5: MODELING ARCHITECTURE & PIPELINE
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Predictive Modeling & Multi-Stage Ensemble",
               "TRIPLE GRADIENT-BOOSTED ENSEMBLE + TARGETED CONDITIONAL REGRESSION CASCADE")
    
    # Left: Architecture Flow Card
    add_card(s5, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.3))
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.6), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Ensemble Architecture & Validation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    arch_points = [
        ("Stratified 5-Fold Cross-Validation", "Athlete-level grouping with preserved class balance across all folds."),
        ("Task A Classifier Ensemble", "Weighted blend: 40% LightGBM + 40% XGBoost + 20% HistGradientBoosting for calibrated probability outputs."),
        ("Task B1: Onset Day Regressor", "Trained on injured athletes, predicting integer offset [1, 30] governed by fatigue accumulation velocity."),
        ("Task B2: Recovery Duration Regressor", "Gradient Boosted Regressor capturing sport-specific anatomical load and injury recovery distributions [5, 20].")
    ]
    for title, desc in arch_points:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_TEXT
        
    # Right: Embed Classification Performance Figure
    if os.path.exists('figures/03_model_classification_performance.png'):
        s5.shapes.add_picture('figures/03_model_classification_performance.png', Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.3))

    # ==========================================
    # SLIDE 6: EVALUATION & RESULTS
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Model Performance & Evaluation Metrics",
               "EVALUATED ACROSS TASK A (F1), TASK B TIMING MAE, AND PENALIZED SKILL SCORES")
    
    # Top 3 Metric Highlight Cards
    metric_cards = [
        ("Task A Classification (F1)", "0.6570", "Precision: 90.6% | AUC: 0.941", BLUE_ACCENT),
        ("Task B1: Onset Hit MAE", "0.82 Days", "Baseline MAE: 7.61 Days", EMERALD),
        ("Task B2: Recovery Hit MAE", "2.90 Days", "Baseline MAE: 3.24 Days", RGBColor(168, 85, 247))
    ]
    for i, (title, val, sub, col) in enumerate(metric_cards):
        left = Inches(0.8 + i * 4.0)
        add_card(s6, left, Inches(1.5), Inches(3.7), Inches(1.6))
        tb = s6.shapes.add_textbox(left + Inches(0.2), Inches(1.6), Inches(3.3), Inches(1.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = MUTED_TEXT
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = col
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = DARK_TEXT
        
    # Embed Timing Accuracy Figure
    if os.path.exists('figures/05_timing_predictions_evaluation.png'):
        s6.shapes.add_picture('figures/05_timing_predictions_evaluation.png', Inches(0.8), Inches(3.4), Inches(11.7), Inches(3.5))

    # ==========================================
    # SLIDE 7: FEATURE IMPORTANCE & INTERPRETABILITY
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Model Interpretability & Biometric Drivers",
               "WHAT SIGNALS WARN OF AN IMPENDING INJURY?")
    
    # Left: Embed Feature Importance Figure
    if os.path.exists('figures/04_feature_importance.png'):
        s7.shapes.add_picture('figures/04_feature_importance.png', Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.3))
        
    # Right: Domain Interpretation Card
    add_card(s7, Inches(7.5), Inches(1.5), Inches(5.0), Inches(5.3))
    tb = s7.shapes.add_textbox(Inches(7.7), Inches(1.7), Inches(4.6), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Biometric Risk Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    drivers = [
        ("Peak Workload Surges", "chronic_TotalSteps_max and step_spike_delta are the #1 and #2 most important features, capturing unmanaged spikes in training intensity."),
        ("Nocturnal Resting Heart Rate", "Elevated minimum heart rate during 2–6 AM indicates chronic autonomic fatigue and incomplete physiological recovery."),
        ("Sleep Regularity & Deficit", "Athletes with irregular sleep patterns (high sleep CV) and cumulative deficits are significantly more vulnerable to soft-tissue failure."),
        ("Sport & Position Load Profile", "High mechanical impact positions (e.g. Football Forwards, Basketball Centers) have distinct fatigue thresholds.")
    ]
    for title, desc in drivers:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_TEXT

    # ==========================================
    # SLIDE 8: CONCLUSION & DEPLOYMENT IMPACT
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. Operational Deployment & Athlete Safety Impact",
               "FROM PREDICTIVE MODELING TO PROACTIVE SPORTS MEDICINE INTERVENTIONS")
    
    # 3 Large Cards
    conclusions = [
        ("Proactive Workload Management", 
         "Enables sports science staff to detect unsafe training spikes (ACWR > 1.5) up to 30 days in advance, allowing for micro-dosing and targeted rest before tissue overload occurs."),
        ("Tailored Rehabilitation Planning",
         "Accurate forecast of recovery duration allows medical teams to schedule staged RTP (Return-to-Play) protocols based on athlete-specific and sport-specific baselines."),
        ("Scalable Edge/Wearable Pipeline",
         "Lightweight gradient boosting models run with sub-second inference latency, ready for deployment directly onto athlete monitoring dashboards and wearable companion apps.")
    ]
    for i, (title, desc) in enumerate(conclusions):
        left = Inches(0.8 + i * 4.0)
        add_card(s8, left, Inches(1.5), Inches(3.7), Inches(5.3))
        tb = s8.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = DARK_TEXT
        
    prs.save(output_pptx)
    print(f"\n>>> Successfully generated presentation deck: {output_pptx}")

if __name__ == '__main__':
    create_deck()
